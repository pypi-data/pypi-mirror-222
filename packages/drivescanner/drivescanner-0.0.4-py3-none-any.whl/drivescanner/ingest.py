from bs4 import UnicodeDammit, BeautifulSoup
import re
from typing import Tuple
from drivescanner.file_index import _replace_backslash
from pypdf import PdfReader
import docx2txt
import openpyxl
import pandas as pd
from pptx import Presentation

non_processable_file_types = ["pyc"]


def _extract_text_from_table(table):
    text = ""
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    text += f"{run.text} "
            text += "\n"
    return text


def _read_pptx(filepath: str):
    try:
        prs = Presentation(filepath)
        content = ""

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            content += f"{run.text} "
                elif shape.has_table:
                    table = shape.table
                    content += _extract_text_from_table(table)
        errormsg = None
    except FileNotFoundError:
        content = None
        errormsg = "File not found"
    except Exception as e:
        content = None
        errormsg = e

    return content, errormsg


def _read_bytes(filepath: str) -> Tuple[bytearray, str]:
    """
    The function reads the contents of a file in binary mode and returns a tuple containing the content
    as a bytearray and an error message if any.

    Args:
      filepath (str): A string representing the file path of the file to be read.

    Returns:
      a tuple containing two values: a bytearray object representing the content of the file at the
    given filepath, and a string representing an error message if an error occurred while reading the
    file (or None if no error occurred).
    """
    try:
        with open(filepath, "rb") as f:
            content = f.read()
            errormsg = None
    except FileNotFoundError:
        content = None
        errormsg = "File not found"
    except Exception as e:
        content = None
        errormsg = e
    return content, errormsg


def _read_docx(filepath: str) -> Tuple[bytearray, str]:
    """
    This function reads the content of a docx file and returns it along with an error message if the
    file is not found.

    Args:
      filepath (str): The filepath parameter is a string that represents the path to a Microsoft Word
    document in the .docx format.

    Returns:
      A tuple containing a bytearray of the processed content of a docx file and an error message (if
    applicable).
    """
    try:
        content = docx2txt.process(filepath)
        errormsg = None
    except FileNotFoundError:
        content = None
        errormsg = "File not found"
    return content, errormsg


def _read_xlsx(filepath: str) -> str:
    try:
        # first count the number of tabs in the Excel
        workbook = openpyxl.load_workbook(filepath)
        sheets = len(workbook.sheetnames)

        # if the Excel contains a single tab, convert file to  df
        if sheets == 1:
            df = pd.read_excel(filepath)

            # drop rows that are completely empty
            df = df.dropna(how="all", axis=0)

            # drop columns that are completely empty
            df = df.dropna(how="all", axis=1)

            # remove NaN values in string
            content = str(df)
            content = content.replace("NaN", "")

        # otherwise, convert the Excel to a dict
        else:
            df = pd.read_excel(filepath, sheet_name=None)

            # remove NaN values in string
            content = str(df)
            content = content.replace("NaN", "")
            errormsg = None

    except FileNotFoundError:
        content = None
        errormsg = "File not found"
    except Exception as e:
        content = None
        errormsg = e
    return content, errormsg


def _decode_bytearray(
    bytestring: bytearray, try_encoding: list[str]
) -> Tuple[str, str]:
    """
    If the bytestring is not None, then try to decode it using the encoding derived
    from the bytestring or try to decode it using the encoding specified in the
    using some predefined encodings

    Args:
      bytestring (bytearray): the bytearray to decode
      try_encoding (list[str]): a list of encodings to try.

    Returns:
      the content of the bytestring as a decoded string
    """
    possible_encoding = UnicodeDammit(bytestring).original_encoding
    content = None
    errormsg = None
    if (possible_encoding is not None) & (possible_encoding not in try_encoding):
        try_encoding.insert(0, possible_encoding)
    for enc in try_encoding:
        if content is None:
            try:
                content = bytestring.decode(enc)
                errormsg = None
            except UnicodeDecodeError:
                errormsg = "File encoding unknown"
            except Exception as e:
                errormsg = str(e)
    return content, errormsg


def _clean_text(text: str) -> str:
    """
    It removes HTML and XML characters, and replaces all new lines and white spaces like tabs with a
    space

    Args:
      text (str): the text to be cleaned

    Returns:
      A string with no new lines or tabs.
    """
    # remove HTML and XML characters
    text = BeautifulSoup(text).get_text()
    # replace all new lines and white spaces like tabs with a space
    text = re.sub(r"\s", " ", text)
    return text


def _pdf_to_string(pdf_file: str) -> Tuple[str, str]:
    """
    It opens the PDF file, reads the text from each page, and returns the text as a string

    Args:
      pdf_file (str): The path to the PDF file.

    Returns:
      A tuple of two strings, the content and/or the error message (if any) while processing the file.
    """
    try:
        with open(pdf_file, "rb") as pdf:
            reader = PdfReader(pdf, strict=False)
            pdf_text = []

            for page in reader.pages:
                content = page.extract_text()
                pdf_text.append(content)

            content = " ".join(pdf_text)
            errormsg = None
    except FileNotFoundError:
        content = None
        errormsg = "File not found"
    except Exception as e:
        content = None
        errormsg = e
    return content, errormsg


def ingest_file(
    filepath: str, extension: str, try_encoding: list[str] = ["utf-8", "iso-8859-1"]
) -> Tuple[str, str]:
    """
    It reads a file, tries to decode it, and then cleans it

    Args:
      filepath (str): str = the path to the file you want to ingest
      try_encoding (list[str]): list[str] = ["utf-8", "iso-8859-1"]

    Returns:
      A string
    """
    filepath = _replace_backslash(filepath)

    if extension in non_processable_file_types:
        # we don't want to process these files
        content = None
        errormsg = f"File of type {extension} is not supported"
    elif extension == "pdf":
        # use PyPDF2
        content, errormsg = _pdf_to_string(filepath)
    elif extension == "docx":
        content, errormsg = _read_docx(filepath)
        if (content is not None) and (errormsg is None):
            content, errormsg = _decode_bytearray(
                bytestring=content, try_encoding=try_encoding
            )
    elif extension in ["pptx"]:
        content, errormsg = _read_pptx(filepath)
    elif extension in ["xlsx"]:
        content, errormsg = _read_xlsx(filepath)
    elif extension in [
        "ipynb",
        "sql",
        "r",
        "json",
        "c",
        "py",
        "xml",
        "html",
        "htm",
        "csv",
        "txt",
    ]:
        # use simple read function
        content, errormsg = _read_bytes(filepath)
        if (content is not None) and (errormsg is None):
            content, errormsg = _decode_bytearray(
                bytestring=content, try_encoding=try_encoding
            )
    else:
        # we can't process these files
        content = None
        errormsg = f"File of type {extension} is not supported"

    if (content is not None) and (errormsg is None):
        content = _clean_text(content)

    return content, errormsg
